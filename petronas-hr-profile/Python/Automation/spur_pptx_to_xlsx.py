from pptx import Presentation
import html
import itertools
import pandas as pd
import re
import win32com.client
from config import *
import save_single_slide


def pptx_to_xlsx(
    ppt_list,
    # xlsx_destination,
    save_slide,
    job_blob_dir
):
    try:
        ppt_list = [x for x in ppt_list if not re.search(r"~\$", x)]

        all_ppt = []
        ur_name_all = []
        source_file = []

        if save_slide == True:
            ppt_instance = win32com.client.Dispatch("PowerPoint.Application")

        for ppt in ppt_list:
            LOGGER.info("[Extract pptx to xlsx] Extracting {}".format(ppt))
            prs = Presentation((ppt))

            ur_name_list = []
            ur_print_list = []
            get_text = []
            bold_list = []
            for idx, slide in enumerate(prs.slides, 1):

                for shape in slide.shapes:
                    if not shape.has_table:
                        continue

                    tbl = shape.table
                    if len(tbl.rows) == 1 or len(tbl.columns) == 1:
                        continue

                    if tbl.cell(1, 0).text.strip() == "" or tbl.cell(1, 1).text.strip() == "":
                        continue
                    if (
                        re.search("Role\s+Purpose", " ".join(tbl.cell(1,0).text.strip().split()), flags=re.I)
                        or re.search("Role\s+Purpose", " ".join(tbl.cell(1,1).text.strip().split()), flags=re.I)
                        # or tbl.cell(0, 1).text.strip() in "Role Purpose"
                        # and tbl.cell(0, 0).text.strip() != "Team"
                    ):
                        source_file.append(re.search(r"[^\\]+$", ppt).group(0))
                        for shape2 in slide.shapes:
                            if not shape2.has_text_frame:
                                continue
                            else:
                                if re.search("\d+\s*-\s*\d+|^UR|^F\d", shape2.text):
                                    title = re.sub("\xa0", " ", shape2.text).strip()
                                    break
                                else:
                                    continue
                    else:
                        continue

                    # slide_index = [idx, idx + 1]
                    if re.search("Page\s*2/2|2/2|continue|cont\.", title, flags=re.I):
                        slide_index = [idx - 1, idx]
                    else:
                        slide_index = [idx]
                    ur_code = (
                        re.search("\d+\s*\-\s*\d+\s*(?=\||\)|\D)|^UR.+?(?=\:)|^F\d.+(?=\|)", title, flags=re.S)
                        .group()
                        .strip()
                    )
                    print(title)
                    print(ur_code)

                    ur_code = re.sub("\s+", "", ur_code)
                    ur_name_list.append(title)
                    ur_print_list.append({title: idx})

                    row_count = len(tbl.rows)
                    col_count = len(tbl.columns)

                    if col_count == 6:
                        start_col = 1
                    elif col_count == 7:
                        start_col = 1
                        col_count = 6
                    else:
                        start_col = 0
                    data = []
                    for r in range(0, row_count):
                        col_text = []
                        for c in range(start_col, col_count):
                            z = ""
                            cell = tbl.cell(r, c)
                            paragraphs = cell.text_frame.paragraphs
                            # if c != start_col:
                            #     z += " | "

                            numbering = 1
                            for paragraph in paragraphs:
                                # if title == "F02-010 | Executive (Brand Strategy)":
                                #     print(paragraph.text.strip())
                                #     print(paragraph._element.xml)
                                text = paragraph.text.strip()

                                if 'type="arabicPeriod"' in paragraph._element.xml and text != "":
                                    text = "{}. ".format(numbering) + text
                                    numbering += 1
                                if 'type="arabicParenR"' in paragraph._element.xml and text != "":
                                    text = "{}) ".format(numbering) + text
                                    numbering += 1
                                if 'type="arabicPlain"' in paragraph._element.xml and text != "":
                                    text = "{} ".format(numbering) + text
                                    numbering += 1

                                if re.search(r"^\d", text):
                                    numbering = int(re.search(r"^\d", text).group(0)) + 1

                                if 'u="sng"' in paragraph._element.xml and text != "":
                                    underline_text_list = re.findall("<a:r>[\s\S]*?<\/a:r>", paragraph._element.xml)

                                    for t in underline_text_list:
                                        if 'u="sng"' in t:
                                            underline_text = html.unescape(
                                                re.search("(?<=<a:t>)[\s\S]+?(?=<\/a:t>)", t).group().strip()
                                            )
                                            if underline_text == "":
                                                continue
                                            underline_text = (
                                                underline_text.replace("(", "\(").replace(")", "\)").replace("*", "\*")
                                            )
                                            if re.search(" \\\\ ", underline_text):
                                                underline_text = underline_text.replace("\\", "\\\\")
                                            underline_text_location = re.search(underline_text, text).span(0)
                                            text = (
                                                text[: underline_text_location[0]]
                                                + "<u>"
                                                + text[underline_text_location[0] : underline_text_location[1]]
                                                + "</u>"
                                                + text[underline_text_location[1] :]
                                            )

                                if "•" in paragraph._element.xml and text != "":
                                    z += "<li>•" + text + "•</li>" + "\n"

                                elif "§" in paragraph._element.xml and text != "":
                                    z += "<li>•" + text + "•</li>" + "\n"

                                elif 'buChar char="o"' in paragraph._element.xml and text != "":
                                    z += "<li>o\t" + text + "o\t</li>" + "\n"

                                elif 'type="alphaLcParenR"' in paragraph._element.xml and text != "":
                                    z += "<li>alphaLcParenR" + text + "alphaLcParenR</li>" + "\n"

                                elif 'type="alphaLcPeriod"' in paragraph._element.xml and text != "":
                                    z += "<li>alphaLcPeriod" + text + "alphaLcPeriod</li>" + "\n"

                                elif 'b="1"' in paragraph._element.xml and text != "":
                                    bold_text_list = re.findall("<a:r>[\s\S]*?<\/a:r>", paragraph._element.xml)
                                    # if title == "F02-010 | Executive (Brand Strategy)":
                                    #     print(bold_text_list)

                                    for t in bold_text_list:
                                        if 'b="1"' in t:
                                            bold_text = html.unescape(
                                                re.search("(?<=<a:t>)[\s\S]+?(?=<\/a:t>)", t).group().strip()
                                            )
                                            # if title == "F02-010 | Executive (Brand Strategy)":
                                            #     print(bold_text)
                                            if bold_text == "":
                                                continue
                                            bold_text = (
                                                bold_text.replace("(", "\(")
                                                .replace(")", "\)")
                                                .replace("*", "\*")
                                                .replace(".", "\.")
                                            )
                                            if re.search(" \\\\ ", bold_text):
                                                bold_text = bold_text.replace("\\", "\\\\")
                                            # print(bold_text)
                                            bold_text_location = re.search(bold_text, text).span(0)
                                            if re.search(r"^\d\W", text):
                                                text = (
                                                    "<strong>"
                                                    + text[: bold_text_location[0]]
                                                    + text[bold_text_location[0] : bold_text_location[1]]
                                                    + "</strong>"
                                                    + text[bold_text_location[1] :]
                                                    + "<br>"
                                                )
                                            else:
                                                text = (
                                                    text[: bold_text_location[0]]
                                                    + "<strong>"
                                                    + text[bold_text_location[0] : bold_text_location[1]]
                                                    + "</strong>"
                                                    + text[bold_text_location[1] :]
                                                    + "<br>"
                                                )
                                            text = re.sub("</strong> +<strong>", "", text)
                                    if re.search(r"^<strong>|^<u><strong>|^\*<strong>", text):
                                        text = "<p>" + text + "</p>"
                                        if re.search("<u><strong>", text):
                                            text = re.sub("<u><strong>", "<p><u><strong>", text)
                                            text = re.sub("</strong></u>", "</strong></u></p>", text)
                                        elif re.search("\*<strong>", text):
                                            text = re.sub("\*<strong>", "<p>\*<strong>", text)
                                        else:
                                            text = re.sub("<strong>", "<p><strong>", text)
                                            text = re.sub("</strong>", "</strong></p>", text)
                                        text = text.strip()
                                        text = re.sub("<br>|</p><p>", "", text)
                                    z += text + "\n"

                                else:
                                    z += text + "<br>" + "\n"
                            col_text.append(z)
                        data.append(col_text)

                    def remove_html_tags(text):
                        """Remove html tags from a string"""
                        import re

                        text = str(text)

                        clean = re.compile("<.*?>")
                        return re.sub(clean, "", text)

                    temp_df = pd.DataFrame(data[2:], columns=data[1])
                    temp_df.columns = [remove_html_tags(x).strip() for x in temp_df.columns]
                    if col_count > 7:
                        final_columns = [
                            x
                            for x in temp_df.columns
                            if re.search(
                                "Role\s*Purpose|Acc|Accountabilities|Challenges|Experience|Key\s*Performance\s*Indicator|KPI",
                                x,
                                flags=re.I | re.S,
                            )
                        ]
                        temp_df = temp_df[final_columns]

                    data_list = []
                    for col in temp_df.columns:
                        # print(col)
                        # if re.search("")
                        data_list.append("<br>".join(temp_df[col].values.tolist()))

                    z = "|".join(data_list)

                    z = re.sub(r"(?<!•</li>)\s+<li>•|^<li>•", "\n<ul>\n<li>", z)
                    z = re.sub(r"•</li>\s+(?!<li>•)", "</li>\n</ul>\n", z)

                    z = re.sub(r"(?<!o\t</li>)\s+<li>o\t", "\n<ul>\n<li>", z)
                    z = re.sub(r"o\t</li>\s+(?!<li>o\t)", "</li>\n</ul>\n", z)

                    ol_format_1 = """<style>
                    ol {
                        counter-reset: list;
                    }
                    ol > li {
                        list-style: none;
                        position: relative;
                    }
                    ol > li:before {
                        counter-increment: list;
                        content: counter(list, lower-alpha) ") ";
                        position: absolute;
                        left: -1.4em;
                    }
                    </style>
                    """

                    z = re.sub(
                        r"(?<!alphaLcParenR</li>)\s+<li>alphaLcParenR",
                        '\n{}\n<ol type="a">\n<li>'.format(ol_format_1),
                        z,
                    )
                    z = re.sub(r"alphaLcParenR</li>\s+(?!<li>alphaLcParenR)", "</li>\n</ol>\n", z)

                    z = re.sub(
                        r"(?<!alphaLcPeriod</li>)\s+<li>alphaLcPeriod",
                        '\n<ol type="a">\n<li>',
                        z,
                    )
                    z = re.sub(r"alphaLcPeriod</li>\s+(?!<li>alphaLcPeriod)", "</li>\n</ol>\n", z)

                    z = re.sub("•|alphaLcParenR|alphaLcPeriod|o\t|\v|\f|\r", "", z)

                    # z = re.sub("</ul>\n<ul>", "<ul>", z)
                    # z = re.sub("</ul>\n<ul>", "<ul>", z)
                    # z = re.sub(r'</ul>\n<br>\n<strong>', '</ul>\n\n<strong>', z)
                    # if title == "F02-010 | Executive (Brand Strategy)":
                    #     print(repr(z))
                    # temp_df = pd.DataFrame()

                    get_text.append(z)

                    if save_slide == True:
                        # Save slide
                        A = save_single_slide.save_single_slide(
                            ppt_instance=ppt_instance,
                            pptx_file_path=ppt,
                            slide_index=slide_index,
                            save_path=job_blob_dir,
                            filename="{}.pptx".format(ur_code)
                        )
                        A.save_slide()
                        LOGGER.info("[Extract pptx to xlsx] Generating {}".format(A.filename.replace("pptx", "pdf")))
                    # print(z)
                # pass

            ur_name_all.append(ur_name_list)
            all_ppt.append(get_text)

            # test = []
            # for key, value in ur_print_list:
            #     if key not in [x for x in ]

        # Stop ppt_instance
        if save_slide == True:
            ppt_instance.Quit()

        # if len(os.listdir(job_blob_dir)) != 0:
        #     for f in os.listdir(job_blob_dir):
        #         shutil.copy2(job_blob_dir + "\\" + f, position_blob_dir)

        text_list = list(itertools.chain.from_iterable(all_ppt))
        ur_name_list = list(itertools.chain.from_iterable(ur_name_all))

        spur_df = (
            pd.DataFrame({"UR": ur_name_list, "text": text_list, "SOURCE": source_file})
            # .drop_duplicates(subset=["UR"])
            .reset_index(drop=True)
        )
        spur_df["UR_CODE"] = spur_df["UR"].apply(
            lambda x: re.search(
                "[A-Za-z]+\d+\s*\-\s*\d+\w*\s*(?=\||\))|\d+\s*\-\s*\d{3}\w*\s*(?=\||:|\-|\)|\w)", x
            ).group()
        )
        spur_df["UR_CODE"] = spur_df["UR_CODE"].apply(lambda x: re.sub("\s+", "", x))
        spur_df["UR_NAME"] = spur_df["UR"].apply(
            lambda x: re.sub(
                "[A-Za-z]+\d+\s*\-\s*\d+\w*\s*(?=\|)|\d+\s*\-\s*\d{3}\w*\s*(?=\||\-|\w)|\||^UR.+?(?=\:)|:|"
                "Do & Deliver|DO,DELIVER|DO , DELIVER|DO , DISPLAY|"
                "DO, DELIVER|DO, DISPLAY|- DO & DELIVER|– DO & DELIVER|- Do & Deliver \(\d/\d\)",
                "",
                x,
            )
            .strip()
            .strip("-")
            .strip("–")
        )
        spur_df["UR_NAME"] = spur_df["UR_NAME"].apply(lambda x: re.sub("Hybrid position\s+\(\)", "", x).strip())
        spur_df[["ROLEPURPOSE", "ACCOUNTABILITIES", "CHALLENGES", "EXPERIENCE", "KPI"]] = spur_df["text"].str.split(
            "\|",
            expand=True,
        )
        spur_df = (
            spur_df[
                [
                    "UR_CODE",
                    "UR_NAME",
                    "ROLEPURPOSE",
                    "ACCOUNTABILITIES",
                    "CHALLENGES",
                    "EXPERIENCE",
                    "KPI",
                    "SOURCE",
                ]
            ]
            .apply(lambda x: x.str.strip())
            # .applymap(lambda x: re.sub("^(?=[upol]+>)", "<", x.strip(" \n<br>")) if re.search("^<br>|<br>$", x) else x)
            # .applymap(lambda x: re.sub("(?<=[<upol/])$", ">", x))
            # .applymap(lambda x: re.sub("\.<br>", "<br>", x))
            # .drop_duplicates(subset=["UR_CODE", "ACCOUNTABILITIES"])
            # .pipe(lambda x: x.groupby("UR_CODE").agg("|".join).reset_index(drop=False))
        ).applymap(lambda x: "<br>".join(list(dict.fromkeys(x.split("|")))) if "|" in x else x)
        spur_df["UR_NAME"] = spur_df["UR_NAME"].apply(
            lambda x: re.sub("–\s+\(Page.+|–\s+\(\d.+|\(\d/\d.+|<br>.+|\s*\d\/\d.*", "", x, flags=re.S).strip().strip("*")
        )
        spur_df = spur_df.replace("–", "-", regex=True)
        # spur_df.to_excel(xlsx_destination, index=False)

        # for i in range(len(spur_df)):
        #     # DESCRIPTION
        #     with open(
        #         job_clob_dir + "\\" + spur_df.iloc[i]["UR_CODE"] + "_DESCRIPTION.txt",
        #         "w",
        #         encoding="utf-8",
        #     ) as f:
        #         f.write(
        #             "<p>"
        #             + spur_df.iloc[i]["ROLEPURPOSE"].strip()
        #             + "</p>"
        #             + "\n\n\n"
        #             + spur_df.iloc[i]["ACCOUNTABILITIES"].strip()
        #         )
        #     # RESPONSIBILITY
        #     with open(
        #         job_clob_dir + "\\" + spur_df.iloc[i]["UR_CODE"] + "_RESPONSIBILITY.txt",
        #         "w",
        #         encoding="utf-8",
        #     ) as f:
        #         f.write(spur_df.iloc[i]["KPI"].strip())
        #     # QUALIFICATION
        #     with open(
        #         job_clob_dir + "\\" + spur_df.iloc[i]["UR_CODE"] + "_QUALIFICATION.txt",
        #         "w",
        #         encoding="utf-8",
        #     ) as f:
        #         f.write(spur_df.iloc[i]["EXPERIENCE"].strip())

        # if (save_slide == True) & (len(os.listdir(job_clob_dir)) != 0):
        #     for f in os.listdir(job_clob_dir):
        #         shutil.copy2(job_clob_dir + "\\" + f, position_clob_dir)

        return spur_df
    except Exception as e:
        raise ValueError(e)